from __future__ import annotations

import re
from collections import Counter
from pathlib import Path
from typing import Iterable

try:
    from docx import Document as DocxDocument
except ImportError:  # pragma: no cover
    DocxDocument = None


ROOT = Path(__file__).resolve().parent
INPUT_DIR = ROOT / "input"
OUTPUT_DIR = ROOT / "output"
PROMPTS_DIR = ROOT / "prompts"
ASSETS_DIR = ROOT / "assets"
PPTX_PATH = OUTPUT_DIR / "讲座整理PPT_初版.pptx"
OUTLINE_PATH = OUTPUT_DIR / "slides_outline.md"
IMAGE_PROMPTS_PATH = PROMPTS_DIR / "image_prompts.md"
MISSING_README = INPUT_DIR / "README_请放入逐字稿.md"
TRANSCRIPT_CANDIDATES = [
    INPUT_DIR / "transcript.docx",
    INPUT_DIR / "transcript.txt",
    INPUT_DIR / "transcript.md",
]

STOPWORDS = set("""
这个 那个 然后 就是 其实 因为 所以 如果 但是 我们 你们 他们 它们 一个 一些 以及 进行 通过 可以 需要 不是 没有 可能 比较 非常 大家 时候 问题 内容 方面 相关 这里 里面 现在 今天 觉得 知道 看到 讲到 提到 这种 这样 那样
""".split())

SECTIONS = [
    ("问题背景", "要先说明问题为什么值得被重视"),
    ("核心判断", "这份材料的中心判断需要被明确表达"),
    ("主要观点", "关键观点要围绕同一主线展开"),
    ("方法步骤", "可执行路径比零散建议更重要"),
    ("案例说明", "用材料中的例子帮助听众理解"),
    ("局限提醒", "边界条件和风险需要提前说明"),
    ("行动建议", "把判断转化为下一步行动"),
    ("总结页", "回到主线并强化结论"),
]


def ensure_dirs() -> None:
    for path in (INPUT_DIR, OUTPUT_DIR, PROMPTS_DIR, ASSETS_DIR):
        path.mkdir(parents=True, exist_ok=True)


def write_missing_readme() -> None:
    MISSING_README.write_text(
        "# 请先放入逐字稿\n\n"
        "当前仓库没有检测到可用于生成 PPT 的逐字稿文件。\n\n"
        "请将逐字稿放入以下任意一个位置后，再运行 `make_ppt.py`：\n\n"
        "1. `input/transcript.docx`\n"
        "2. `input/transcript.txt`\n"
        "3. `input/transcript.md`\n\n"
        "脚本会按以上顺序优先读取文件，并根据逐字稿内容生成 PPT、提纲和配图提示词。\n",
        encoding="utf-8",
    )


def find_transcript() -> Path | None:
    return next((p for p in TRANSCRIPT_CANDIDATES if p.exists()), None)


def read_transcript(path: Path) -> str:
    if path.suffix.lower() == ".docx":
        if DocxDocument is None:
            raise RuntimeError("读取 docx 需要安装 python-docx：pip install -r requirements.txt")
        doc = DocxDocument(str(path))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    return path.read_text(encoding="utf-8-sig")


def clean_text(text: str) -> str:
    text = re.sub(r"\s+", " ", text)
    text = re.sub(r"(嗯|啊|呃|额|这个|那个)[，,、\s]*", "", text)
    return text.strip()


def split_sentences(text: str) -> list[str]:
    parts = re.split(r"(?<=[。！？；!?;])\s*", text)
    return [p.strip(" ，,。；;：:") for p in parts if len(p.strip()) >= 8]


def keywords(sentences: Iterable[str], topn: int = 10) -> list[str]:
    words: list[str] = []
    for sentence in sentences:
        words.extend(re.findall(r"[\u4e00-\u9fff]{2,6}", sentence))
    counts = Counter(w for w in words if w not in STOPWORDS)
    return [w for w, _ in counts.most_common(topn)]


def pick_points(sentences: list[str], start: int, count: int = 4) -> list[str]:
    if not sentences:
        return ["逐字稿内容较短，建议补充更多事实、案例和论证材料"]
    points = []
    idx = start
    while len(points) < count and idx < len(sentences):
        item = sentences[idx]
        item = re.sub(r"^(所以|然后|但是|因为|那么|首先|其次|最后)[，,、\s]*", "", item)
        if len(item) > 46:
            item = item[:44] + "……"
        if item and item not in points:
            points.append(item)
        idx += 1
    return points or ["围绕本页主题提炼逐字稿中的关键表达"]


def build_slides(text: str) -> list[dict]:
    cleaned = clean_text(text)
    sentences = split_sentences(cleaned)
    kws = keywords(sentences)
    theme = "、".join(kws[:3]) if kws else "讲座核心内容"
    target_count = 15 if len(sentences) < 80 else min(20, max(15, len(sentences) // 12))

    slides = [
        {"title": f"围绕{theme}建立清晰共识", "points": ["基于逐字稿整理", "适用于汇报、讲座或课程交流", "突出结构、判断与行动建议"], "notes": "开场说明本 PPT 根据逐字稿整理，重点不是逐句复述，而是帮助听众把握主线。"},
        {"title": "本次分享将按问题、判断、方法和行动展开", "points": [name for name, _ in SECTIONS], "notes": "提示听众接下来会先讲背景，再讲判断与观点，最后落到行动建议。"},
    ]
    remaining = target_count - 3
    per_section = max(1, remaining // len(SECTIONS))
    cursor = 0
    for name, judgement in SECTIONS:
        for i in range(per_section):
            if len(slides) >= target_count - 1:
                break
            points = pick_points(sentences, cursor, 4)
            cursor += 4
            slides.append({"title": judgement if i == 0 else f"{name}需要进一步转化为可理解的表达", "points": points[:5], "notes": f"本页围绕“{name}”展开，讲解时应结合逐字稿原意，删去口语重复，突出与听众相关的含义。"})
    slides.append({"title": "总结应回到核心判断并推动后续行动", "points": ["再次确认问题背景", "重申最重要的判断", "保留必要的风险提醒", "给出可立即执行的下一步"], "notes": "收束全篇，帮助听众记住结论，并明确会后可以做什么。"})
    return slides[:20]


def add_notes(slide, text: str) -> None:
    notes = slide.notes_slide.notes_text_frame
    notes.text = text


def add_textbox(slide, left, top, width, height, text, size=24, bold=False, color=None):
    from pptx.dml.color import RGBColor
    from pptx.util import Pt
    if color is None:
        color = RGBColor(31, 49, 68)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def create_ppt(slides_data: list[dict]) -> None:
    from pptx import Presentation
    from pptx.util import Inches
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for index, data in enumerate(slides_data, 1):
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(248, 250, 252)
        if index == 1:
            add_textbox(slide, Inches(0.9), Inches(2.0), Inches(11.5), Inches(0.8), data["title"], 34, True)
            add_textbox(slide, Inches(0.95), Inches(3.0), Inches(10.5), Inches(0.6), "根据逐字稿结构化整理", 20, False, RGBColor(80, 96, 112))
        else:
            add_textbox(slide, Inches(0.65), Inches(0.35), Inches(11.8), Inches(0.6), data["title"], 28, True)
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(1.1), Inches(12.0), Inches(0.02))
            line.fill.solid(); line.fill.fore_color.rgb = RGBColor(54, 96, 146); line.line.fill.background()
            top = 1.45
            for n, point in enumerate(data["points"][:5], 1):
                add_textbox(slide, Inches(1.0), Inches(top), Inches(0.45), Inches(0.35), f"{n}.", 19, True, RGBColor(54, 96, 146))
                add_textbox(slide, Inches(1.5), Inches(top), Inches(10.8), Inches(0.45), point, 20)
                top += 0.85
        add_textbox(slide, Inches(11.8), Inches(7.0), Inches(0.8), Inches(0.2), f"{index:02d}", 10, False, RGBColor(120, 130, 140))
        add_notes(slide, data["notes"])
    prs.save(PPTX_PATH)


def write_outline(slides_data: list[dict]) -> None:
    lines = ["# 讲座整理 PPT 提纲\n"]
    for i, s in enumerate(slides_data, 1):
        lines.append(f"## 第 {i:02d} 页：{s['title']}\n")
        lines.extend(f"- {p}" for p in s["points"])
        lines.append(f"\n**演讲者备注：** {s['notes']}\n")
    OUTLINE_PATH.write_text("\n".join(lines), encoding="utf-8")


def write_image_prompts(slides_data: list[dict]) -> None:
    lines = ["# PPT 配图提示词\n", "第一轮不直接生成图片；以下提示词供后续使用 image2。\n"]
    for i, s in enumerate(slides_data, 1):
        name = f"slide_{i:02d}.png"
        lines.append(f"## 第 {i:02d} 页：{s['title']}")
        lines.append("- 建议图片类型：稳重的抽象背景、真实会议/课堂场景或简洁信息图")
        lines.append(f"- 中文图片提示词：为中文汇报型 PPT 生成一张配图，主题为“{s['title']}”，画面清晰、稳重、留白充足，不要出现夸张表情和广告元素。")
        lines.append(f"- English prompt: Create a calm, professional presentation image for the slide titled '{s['title']}', clear composition, subtle colors, enough negative space, no cartoonish or advertising style.")
        lines.append("- 图片风格建议：蓝灰色系、简洁、商务汇报风、适合 16:9 PPT")
        lines.append(f"- 图片文件建议命名：`{name}`\n")
    IMAGE_PROMPTS_PATH.write_text("\n".join(lines), encoding="utf-8")


def main() -> int:
    ensure_dirs()
    transcript = find_transcript()
    if transcript is None:
        write_missing_readme()
        print(f"未找到逐字稿，已创建 {MISSING_README.relative_to(ROOT)}。")
        return 0
    text = read_transcript(transcript)
    if not clean_text(text):
        raise RuntimeError(f"逐字稿为空：{transcript}")
    slides = build_slides(text)
    write_outline(slides)
    write_image_prompts(slides)
    create_ppt(slides)
    print(f"已生成 {PPTX_PATH.relative_to(ROOT)}，共 {len(slides)} 页。")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
